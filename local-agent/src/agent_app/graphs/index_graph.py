from __future__ import annotations
from typing_extensions import TypedDict
from typing import List, Dict, Any
import os, pathlib, mimetypes
from datetime import datetime

from langgraph.graph import StateGraph, START, END

from agent_app.db import init_db, SessionLocal
from agent_app.utils import sha256_bytes, is_text_like  # is_text_like no longer used here (kept for compat)
from agent_app.chunking import chunk_text
from agent_app.embedding import Embedder
from agent_app.vectorstore import get_collection
from agent_app.config import SETTINGS

# ---- Optional imports for parsers (fail gracefully if missing) ----
try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except Exception:
    PptxPresentation = None


# ---- Discovery rules ----
ALLOW_EXT = {
    ".txt", ".md", ".rst", ".py", ".java",
    ".pdf", ".docx", ".pptx",
}
EXCLUDE_DIRS = {
    ".git", ".venv", "__pycache__", "node_modules",
    "data", "chroma", "checkpoints",
}

# Fallback MIME map for Office formats if mimetypes misses them
OFFICE_MIME = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class IndexState(TypedDict, total=False):
    mode: str                 # "full" | "incremental"
    roots: List[str]
    model: str
    force_reembed: bool
    files: List[Dict[str, Any]]
    changed: List[Dict[str, Any]]
    chunks: List[Dict[str, Any]]
    embeddings: List[List[float]]
    stats: Dict[str, Any]
    errors: List[str]


def _guess_mime(path: str, ext: str) -> str:
    return mimetypes.guess_type(path)[0] or OFFICE_MIME.get(ext, "application/octet-stream")


def discover(s: IndexState) -> IndexState:
    files: list[dict] = []
    for root in s.get("roots", []):
        rp = pathlib.Path(root)
        if not rp.exists():
            continue
        for dirpath, dirnames, filenames in os.walk(rp):
            # prune noisy dirs in-place
            dirnames[:] = [d for d in dirnames if d not in EXCLUDE_DIRS]
            for fn in filenames:
                p = pathlib.Path(dirpath) / fn
                if not p.is_file():
                    continue
                ext = p.suffix.lower()
                if ext not in ALLOW_EXT:
                    continue

                # stat + robust sha256 (content if possible; fallback to size+mtime)
                st = p.stat()
                try:
                    data = p.read_bytes()
                    sha = sha256_bytes(data)
                except Exception:
                    sha = f"{ext or 'file'}:{st.st_size}:{st.st_mtime_ns}"

                files.append({
                    "path": str(p),
                    "bytes": st.st_size,
                    "mtime_ns": st.st_mtime_ns,
                    "sha256": sha,
                    "mime": _guess_mime(str(p), ext),
                    "ext": ext,
                })

    s["files"] = files
    s.setdefault("stats", {})["discovered"] = len(files)
    return s


def diff(s: IndexState) -> IndexState:
    """Decide which files to (re)embed."""
    from sqlalchemy import text
    # Full re-embed if explicitly requested
    if s.get("force_reembed"):
        s["changed"] = list(s.get("files", []))
        s.setdefault("stats", {})["changed"] = len(s["changed"])
        return s

    changed: list[dict] = []
    mode = s.get("mode", "incremental")
    with SessionLocal() as db:
        known = {row[0]: row[1] for row in db.execute(text("SELECT path, sha256 FROM files"))}
    for f in s["files"]:
        if mode == "full" or f["path"] not in known or known[f["path"]] != f["sha256"]:
            changed.append(f)
    s["changed"] = changed
    s.setdefault("stats", {})["changed"] = len(changed)
    return s


def _read_text_any(path: str) -> tuple[str, str, str]:
    """
    Returns (text, mime, ext). ext is lowercase like '.pdf'.
    Raises a readable error if a needed parser is missing.
    """
    p = pathlib.Path(path)
    ext = p.suffix.lower()

    if ext in {".txt", ".md", ".rst", ".py", ".java"}:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(), "text/plain", ext
        except Exception:
            return "", "text/plain", ext

    if ext == ".pdf":
        if PdfReader is None:
            raise RuntimeError("PDF support requires 'pypdf'. Install it to parse PDFs.")
        try:
            text_parts: list[str] = []
            with open(path, "rb") as f:
                pdf = PdfReader(f)
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    text_parts.append(t)
            return "\n".join(text_parts), "application/pdf", ext
        except Exception:
            return "", "application/pdf", ext

    if ext == ".docx":
        if DocxDocument is None:
            raise RuntimeError("DOCX support requires 'python-docx'. Install it to parse Word files.")
        try:
            doc = DocxDocument(path)
            parts: list[str] = [p.text for p in doc.paragraphs]
            # include table text
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            return "\n".join(parts), OFFICE_MIME[".docx"], ext
        except Exception:
            return "", OFFICE_MIME[".docx"], ext

    if ext == ".pptx":
        if PptxPresentation is None:
            raise RuntimeError("PPTX support requires 'python-pptx'. Install it to parse PowerPoint files.")
        try:
            prs = PptxPresentation(path)
            parts: list[str] = []
            for slide in prs.slides:
                # text from shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                # speaker notes, if present
                if getattr(slide, "has_notes_slide", False):
                    notes = getattr(slide.notes_slide, "notes_text_frame", None)
                    if notes and getattr(notes, "text", None):
                        parts.append(notes.text)
            return "\n".join(parts), OFFICE_MIME[".pptx"], ext
        except Exception:
            return "", OFFICE_MIME[".pptx"], ext

    # unknown -> skip
    return "", "application/octet-stream", ext


def parse_chunk(s: IndexState) -> IndexState:
    texts: list[dict] = []
    for f in s.get("changed", []):
        try:
            data, mime, ext = _read_text_any(f["path"])
        except RuntimeError as e:
            # Missing parser: record a clear error, skip file gracefully
            s.setdefault("errors", []).append(f"parse-missing-parser:{f['path']}:{e}")
            continue

        if not data or not data.strip():
            s.setdefault("errors", []).append(f"parse-empty:{f['path']}")
            continue

        # carry through updated mime/ext in case discover guess differed
        f["mime"] = mime
        f["ext"] = ext
        texts.append({"file": f, "text": data})

    s["chunks"] = []
    for t in texts:
        spans = chunk_text(t["text"], target_tokens=800, overlap=80)
        for idx, (cstart, cend, chunk) in enumerate(spans):
            s["chunks"].append({
                "file": t["file"],
                "chunk_idx": idx,
                "char_start": cstart,
                "char_end": cend,
                "text": chunk
            })
    s.setdefault("stats", {})["chunks"] = len(s["chunks"])
    return s


async def embed_batch(s: IndexState) -> IndexState:
    if not s.get("chunks"):
        s["embeddings"] = []
        return s
    embedder = Embedder(model=s.get("model"))
    vecs: list[list[float]] = []
    B = 64
    texts = [c["text"] for c in s["chunks"]]
    for i in range(0, len(texts), B):
        vecs.extend(await embedder.embed_texts(texts[i:i+B]))
    s["embeddings"] = vecs
    return s


def upsert_vectors(s: IndexState) -> IndexState:
    col = get_collection()
    ids, docs, metas = [], [], []
    for i, c in enumerate(s.get("chunks", [])):
        fileinfo = c["file"]
        vid = f"{fileinfo['sha256']}:{c['chunk_idx']}"
        ids.append(vid)
        docs.append(c["text"])
        metas.append({
            "path": fileinfo["path"],
            "sha256": fileinfo["sha256"],
            "chunk_idx": c["chunk_idx"],
            "mime": fileinfo.get("mime"),
            "ext": fileinfo.get("ext"),                 # <-- NEW: extension metadata for better filters
            "mtime_ns": fileinfo.get("mtime_ns"),       # <-- useful for recency filters
            "embedding_model": s.get("model") or SETTINGS.embedding_model,
        })
    if ids:
        col.upsert(ids=ids, documents=docs, metadatas=metas, embeddings=s.get("embeddings", None))
    s.setdefault("stats", {})["upserted"] = len(ids)
    return s


def commit_run(s: IndexState) -> IndexState:
    from sqlalchemy import text
    now = datetime.utcnow().isoformat()
    with SessionLocal() as db:
        for f in s.get("changed", []):
            db.execute(text("""
            INSERT INTO files(path, bytes, mtime_ns, sha256, mime, last_indexed_at)
            VALUES(:path,:bytes,:mtime_ns,:sha256,:mime,:now)
            ON CONFLICT(path) DO UPDATE SET
              bytes=excluded.bytes, mtime_ns=excluded.mtime_ns, sha256=excluded.sha256,
              mime=excluded.mime, last_indexed_at=excluded.last_indexed_at
            """), {**f, "now": now})
        db.commit()
    return s


def build_graph(*, checkpointer=None):
    """Compile the graph with an optional checkpointer (injected by FastAPI at runtime)."""
    init_db()
    g = StateGraph(IndexState)
    g.add_node("discover", discover)
    g.add_node("diff", diff)
    g.add_node("parse_chunk", parse_chunk)
    g.add_node("embed_batch", embed_batch)
    g.add_node("upsert_vectors", upsert_vectors)
    g.add_node("commit_run", commit_run)

    g.add_edge(START, "discover")
    g.add_edge("discover", "diff")
    g.add_edge("diff", "parse_chunk")
    g.add_edge("parse_chunk", "embed_batch")
    g.add_edge("embed_batch", "upsert_vectors")
    g.add_edge("upsert_vectors", "commit_run")
    g.add_edge("commit_run", END)

    return g.compile(checkpointer=checkpointer)


# Expose only NO-checkpointer graph at import (for Studio)
index_graph_api = build_graph(checkpointer=None)
